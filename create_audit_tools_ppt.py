from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "Les Outils et Techniques de l'Audit Interne"
subtitle = slide.placeholders[1]
subtitle.text = "Une Approche Méthodologique Complète"

# Slide 2: Plan
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Plan de Présentation"
content = slide.placeholders[1].text_frame
content.text = "I. Introduction aux Outils d'Audit\nII. Outils d'Interrogation\nIII. Outils de Description\nIV. Outils d'Analyse et d'Évaluation\nV. Outils Informatiques et Numériques\nVI. Conclusion"

# Slide 3: Introduction
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "I. Introduction aux Outils d'Audit"
content = slide.placeholders[1].text_frame
content.text = "Les outils d'audit interne sont des instruments méthodologiques permettant de collecter des informations, analyser les processus et évaluer l'efficacité du contrôle interne.\n\nIls suivent une approche structurée et systématique pour assurer la fiabilité des conclusions."

# Slide 4: Outils d'Interrogation
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "II. Outils d'Interrogation"
content = slide.placeholders[1].text_frame
content.text = "1. Les Entretiens/Interviews\n   - Préparés et structurés\n   - Ouverts ou directifs\n   - Avec enregistrement et compte-rendu\n\n2. Les Questionnaires\n   - Fermés (oui/non) ou ouverts\n   - Auto-administrés ou guidés\n   - Pour collecter des données standardisées"

# Slide 5: Outils de Description
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "III. Outils de Description"
content = slide.placeholders[1].text_frame
content.text = "1. L'Organigramme Fonctionnel\n   - Représentation des relations hiérarchiques\n   - Flux d'information et responsabilités\n\n2. Les Flowcharts/Organigrammes de Processus\n   - Description séquentielle des opérations\n   - Identification des points de contrôle\n\n3. Les Tableaux de Bord\n   - Indicateurs clés de performance\n   - Suivi des objectifs opérationnels"

# Slide 6: Outils d'Analyse
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "IV. Outils d'Analyse et d'Évaluation"
content = slide.placeholders[1].text_frame
content.text = "1. L'Analyse des Risques\n   - Matrice de risques (probabilité x impact)\n   - Cartographie des zones à risque\n\n2. Les Tests de Contrôle\n   - Tests sur pièces (vérification documentaire)\n   - Tests sur le terrain (observation directe)\n\n3. L'Échantillonnage Statistique\n   - Sondage aléatoire ou stratifié\n   - Détermination de la taille d'échantillon"

# Slide 7: Outils Numériques
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "V. Outils Informatiques et Numériques"
content = slide.placeholders[1].text_frame
content.text = "1. L'Analyse de Données (Data Analytics)\n   - Traitement de grands volumes de données\n   - Détection de tendances et anomalies\n\n2. Les Outils d'Audit Assisté par Ordinateur (CAAT)\n   - ACL, IDEA pour l'analyse automatisée\n   - Vérification des contrôles informatiques\n\n3. Les Logiciels de Gestion des Risques\n   - Intégration avec ERP/SAP\n   - Automatisation des contrôles"

# Slide 8: Démarche Méthodologique
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "VI. Démarche Méthodologique d'Utilisation"
content = slide.placeholders[1].text_frame
content.text = "1. Phase de Préparation\n   - Sélection des outils appropriés\n   - Définition des objectifs\n\n2. Phase de Collecte\n   - Application des outils choisis\n   - Documentation systématique\n\n3. Phase d'Analyse\n   - Croisement des informations\n   - Évaluation des contrôles\n\n4. Phase de Conclusion\n   - Rédaction des recommandations\n   - Suivi des actions correctives"

# Slide 9: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Conclusion"
content = slide.placeholders[1].text_frame
content.text = "Les outils d'audit interne constituent un arsenal méthodologique essentiel pour assurer l'efficacité du contrôle interne.\n\nLeur utilisation combinée et adaptée au contexte permet d'atteindre les objectifs d'assurance et de conseil de l'audit interne.\n\nL'évolution technologique offre de nouveaux outils puissants pour renforcer la valeur ajoutée de l'audit."

# Slide 10: Références
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Références"
content = slide.placeholders[1].text_frame
content.text = "• Normes Internationales d'Audit Interne (IIA)\n• Cadre de Référence des Pratiques Professionnelles\n• Ouvrages sur le contrôle interne et l'audit\n• Guides méthodologiques d'audit interne\n• Logiciels spécialisés en audit (ACL, IDEA)"

# Slide 11: Merci
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "Merci de votre attention"
subtitle = slide.placeholders[1]
subtitle.text = "Questions ?"

# Save the presentation
prs.save('/vercel/sandbox/outils_audit_interne_presentation.pptx')
print("Présentation PowerPoint créée avec succès : outils_audit_interne_presentation.pptx")