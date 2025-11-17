from pptx import Presentation
import json

# Load the presentation
prs = Presentation('/vercel/sandbox/uploads/Audit Interne Rappeel .pptx')

# Extract content
slides_content = []

for i, slide in enumerate(prs.slides):
    slide_data = {
        'slide_number': i + 1,
        'texts': [],
        'layout': slide.slide_layout.name
    }
    
    # Extract all text from shapes
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            if shape.text.strip():
                slide_data['texts'].append(shape.text.strip())
        
        # Check for tables
        if shape.has_table:
            table_data = []
            for row in shape.table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                table_data.append(row_data)
            slide_data['table'] = table_data
    
    slides_content.append(slide_data)

# Print the extracted content
print(json.dumps(slides_content, indent=2, ensure_ascii=False))
