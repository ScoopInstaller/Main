#!/usr/bin/env python3
from pptx import Presentation
import sys

def analyze_presentation(file_path):
    prs = Presentation(file_path)

    print(f"\n=== PRESENTATION ANALYSIS ===")
    print(f"Total slides: {len(prs.slides)}")
    print(f"Slide dimensions: {prs.slide_width} x {prs.slide_height}")

    for idx, slide in enumerate(prs.slides, start=1):
        print(f"\n--- SLIDE {idx} ---")
        print(f"Layout: {slide.slide_layout.name}")

        # Extract text from shapes
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content.append(shape.text.strip())

        if text_content:
            print("Text content:")
            for text in text_content:
                print(f"  - {text}")

        # Check for tables
        tables = [s for s in slide.shapes if s.shape_type == 19]  # MSO_SHAPE_TYPE.TABLE
        if tables:
            print(f"  Contains {len(tables)} table(s)")

        # Check for images
        images = [s for s in slide.shapes if s.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE
        if images:
            print(f"  Contains {len(images)} image(s)")

if __name__ == "__main__":
    file_path = "/vercel/sandbox/uploads/Audit Interne Rappeel .pptx"
    analyze_presentation(file_path)
